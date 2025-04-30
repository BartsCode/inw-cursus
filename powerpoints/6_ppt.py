import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Configuration ---
OUTPUT_DIR = 'powerpoints'
OUTPUT_FILENAME = os.path.join(OUTPUT_DIR, 'Hoofdstuk_6.pptx')

# Define some basic colors for syntax highlighting
COLOR_KEYWORD = RGBColor(0x00, 0x00, 0xFF) # Blue
COLOR_STRING = RGBColor(0xFF, 0x00, 0x00)  # Red
COLOR_COMMENT = RGBColor(0x00, 0x80, 0x00) # Green
COLOR_DEFAULT = RGBColor(0x00, 0x00, 0x00) # Black

PYTHON_KEYWORDS = [
    'def', 'if', 'else', 'elif', 'for', 'in', 'while', 'return', 'True', 'False', 'None', 'not', 'and', 'or', 'break', 'continue', 'range', 'len', 'print', 'set'
]

# --- Helper Function for Code Formatting ---
def add_formatted_code(tf, code_lines):
    """Adds code lines to a text frame with basic syntax highlighting."""
    for line in code_lines:
        p = tf.add_paragraph()
        p.font.name = 'Courier New' # Use a monospace font
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

        # Basic highlighting logic
        stripped_line = line.strip()
        if stripped_line.startswith('#'):
            # Whole line is a comment
            run = p.add_run()
            # Preserve original indentation for comments
            run.text = line 
            run.font.color.rgb = COLOR_COMMENT
        else:
            # Handle potential indentation
            indentation = line[:len(line) - len(line.lstrip())]
            if indentation:
                run = p.add_run()
                run.text = indentation
                run.font.color.rgb = COLOR_DEFAULT # Keep default color for space
            
            # Split line by spaces, filtering out empty strings from multiple spaces
            words = [word for word in stripped_line.split(' ') if word]
            
            for i, word in enumerate(words):
                # Add space back if not the first word
                if i > 0:
                    run = p.add_run()
                    run.text = ' ' 
                
                run = p.add_run()
                run.text = word
                
                # Check for keywords
                if word in PYTHON_KEYWORDS:
                    run.font.color.rgb = COLOR_KEYWORD
                # Basic check for strings 
                elif (word.startswith('\"') and word.endswith('\"')) or \
                     (word.startswith('\'') and word.endswith('\'')):
                    run.font.color.rgb = COLOR_STRING
                # Basic check for comments at end of line (handle carefully)
                elif '#' in word and not ((word.startswith('\"') and word.endswith('\"')) or (word.startswith('\'') and word.endswith('\''))):
                     try:
                         comment_start = word.index('#')
                         # Part before comment
                         run.text = word[:comment_start]
                         run.font.color.rgb = COLOR_DEFAULT # Or determine color based on word part
                         
                         # Comment part
                         run_comment = p.add_run()
                         run_comment.text = word[comment_start:]
                         run_comment.font.color.rgb = COLOR_COMMENT
                     except ValueError: # Should not happen if '#' is in word, but be safe
                         run.font.color.rgb = COLOR_DEFAULT
                else:
                    run.font.color.rgb = COLOR_DEFAULT

# --- Main Presentation Function ---
def create_chapter_6_presentation():
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1] # Title and Content
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    blank_slide_layout = prs.slide_layouts[6] # Blank

    # == Slide 1: Title Slide ==
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hoofdstuk 6\nIntroductie tot Algoritmen en Efficiëntie"
    slide.placeholders[1].text = "Concepten, meten en analyseren"

    # == Section 6.1: Wat is een Algoritme? ==
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "6.1 Wat is een Algoritme?"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Een eindige reeks van ondubbelzinnige instructies om een probleem op te lossen."
    p = tf.add_paragraph()
    p.text = "Eigenschappen:"
    p.level = 1
    points = ["Correctheid", "Eindigheid", "Ondubbelzinnigheid", "Input/Output", "Efficiëntie"]
    for point in points:
        p = tf.add_paragraph(); p.text = point; p.level = 2

    # == Section 6.2: Efficiëntie: Waarom Belangrijk? ==
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "6.2 Efficiëntie: Waarom Belangrijk?"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Middelen: Tijd (operaties) & Ruimte (geheugen)."
    p = tf.add_paragraph(); p.text = "Belangrijk voor:"; p.level = 1
    points = ["Gebruikerservaring (snelheid)", "Schaalbaarheid (grote data)", "Kosten (cloud, servers)", "Beperkte apparaten", "Duurzaamheid (energie)"]
    for point in points:
        p = tf.add_paragraph(); p.text = point; p.level = 2

    # == Section 6.3: Big O Notatie ==
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "6.3 Meten van Efficiëntie: Big O Notatie"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Beschrijft hoe Tijd/Ruimte groeit als input (`n`) groeit."
    p = tf.add_paragraph(); p.text = "Focus op groeisnelheid (orde van grootte)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Negeert constanten en lagere-orde termen."; p.level = 1
    p = tf.add_paragraph(); p.text = "Geeft bovengrens (worst-case)."; p.level = 1

    # Common Big O Complexities (Slide 1)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Veelvoorkomende Big O Complexiteiten (1)"
    tf = shapes.placeholders[1].text_frame
    tf.text = "O(1): Constante tijd (bv. `lijst[i]`) - Vlakke lijn"
    p = tf.add_paragraph(); p.text = "O(log n): Logaritmisch (bv. binair zoeken) - Vlakt snel af"
    p = tf.add_paragraph(); p.text = "O(n): Lineair (bv. simpele for-loop) - Rechte stijgende lijn"

    # Common Big O Complexities (Slide 2)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Veelvoorkomende Big O Complexiteiten (2)"
    tf = shapes.placeholders[1].text_frame
    tf.text = "O(n log n): Lineairithmisch (bv. efficiënt sorteren) - Stijgt iets sneller dan O(n)"
    p = tf.add_paragraph(); p.text = "O(n²): Kwadratisch (bv. geneste loop n*n) - Steeds steiler (parabool)"
    p = tf.add_paragraph(); p.text = "O(2^n): Exponentieel (bv. brute-force combinaties) - Explosieve groei"

    # Open Question - Big O Comparison
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Denkvraag: Impact van Groei"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Stel algoritme A is O(n²) en B is O(n)."
    p = tf.add_paragraph(); p.text = "Voor n=10 doet A 100 operaties, B doet 10."; p.level = 1
    p = tf.add_paragraph(); p.text = "Voor n=1.000.000 doet A 1 biljoen (10¹²) operaties, B doet 1 miljoen (10⁶)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Waarom is dit verschil in complexiteit (O(n²) vs O(n)) zo cruciaal voor grote datasets?"
    p = tf.add_paragraph(); p.text = "Wat zijn de praktische gevolgen?"

    # == Section 6.4: Analyse van Eenvoudige Algoritmen ==
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "6.4 Analyse van Eenvoudige Algoritmen"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Basisregels:"
    p = tf.add_paragraph(); p.text = "Enkele statements = O(1)"; p.level = 1
    p = tf.add_paragraph(); p.text = "Sequentie = Hoogste O(...)"; p.level = 1
    p = tf.add_paragraph(); p.text = "If/Else = Conditie + Duurste tak"; p.level = 1
    p = tf.add_paragraph(); p.text = "Loop = Aantal Iteraties * Complexiteit Body"; p.level = 1
    p = tf.add_paragraph(); p.text = "Geneste Loops = Product complexiteiten"; p.level = 1

    # Analysis of List Operations
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Complexiteit Lijstoperaties (lengte n)"
    tf = shapes.placeholders[1].text_frame
    tf.text = "O(1): `lijst[i]`, `append()`, `pop()` (einde), `len()`"
    p = tf.add_paragraph(); p.text = "O(n): `insert()`, `pop(i)`, `x in lijst` (zoeken)"
    p = tf.add_paragraph(); p.text = "O(n log n): `sort()`, `sorted()`"

    # Open Question - Simple Analysis
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Denkvraag: Analyseer Deze Code"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Wat is de Big O tijdcomplexiteit van deze functie?"
    
    code_snippet = [
        "def bereken_som(getallen_lijst):",
        "  # getallen_lijst heeft lengte n",
        "  som = 0             # O(?)",
        "  for getal in getallen_lijst: # O(?)",
        "    som += getal      # O(?)",
        "  return som          # O(?)" 
    ]
    add_formatted_code(tf, code_snippet)
    p = tf.add_paragraph(); p.text = "Hoe schaalt de uitvoeringstijd als de lijst 10x zo lang wordt?"

    # == Section 6.5: Oefeningen ==
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "6.5 Oefeningen: Analyse & Vergelijking"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Laten we enkele voorbeelden uit de oefeningen bekijken..."

    # Oefening 1C - Nested Loop
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Denkvraag: Geneste Loop"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Wat is de Big O tijdcomplexiteit?"
    code_snippet = [
        "def snippet_c(n):",
        "  for i in range(n):      # Hoe vaak?",
        "    for j in range(n):    # Hoe vaak?",
        "      print(i * j)        # O(?)"
    ]
    add_formatted_code(tf, code_snippet)

    # Oefening 3 - Duplicates Comparison
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Denkvraag: Duplicaten Vinden"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Twee methoden om duplicaten te vinden. Welke is sneller voor grote lijsten?"
    
    p = tf.add_paragraph(); p.text = "Methode 1: Geneste lussen (Vergelijk elk paar)"; p.level=1
    p = tf.add_paragraph(); p.text = "Methode 2: Gebruik een `set` (Check `in` voor elk element)"; p.level=1
    
    p = tf.add_paragraph(); p.text = "Denk na over de Big O van elke methode voor je naar de code kijkt."

    # == End Slide ==
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    shapes = slide.shapes
    shapes.title.text = "Vragen?"


    # --- Save Presentation ---
    try:
        if not os.path.exists(OUTPUT_DIR):
            os.makedirs(OUTPUT_DIR)
        prs.save(OUTPUT_FILENAME)
        print(f"Presentatie '{OUTPUT_FILENAME}' is succesvol gegenereerd.")
    except Exception as e:
        print(f"Fout bij opslaan presentatie: {e}")

# --- Run the function ---
if __name__ == "__main__":
    create_chapter_6_presentation()
