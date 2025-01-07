import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_chapter_3_1_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hoofdstuk 3.1\nVergelijkingen en Voorwaarden"
    slide.placeholders[1].text = "Leer hoe je vergelijkingen en voorwaarden gebruikt in Python"

    # Slide 2: Vergelijkingsoperatoren
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Vergelijkingsoperatoren"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    operators = [
        ("`==`", "Is gelijk aan"),
        ("`!=`", "Is niet gelijk aan"),
        ("`>`", "Is groter dan"),
        ("`<`", "Is kleiner dan"),
        ("`>=`", "Is groter dan of gelijk aan"),
        ("`<=`", "Is kleiner dan of gelijk aan")
    ]

    for symbol, meaning in operators:
        p = tf.add_paragraph()
        p.text = f"{symbol} : {meaning}"
        p.level = 1

    # Slide 3: Voorbeelden van Vergelijkingen
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Voorbeelden van Vergelijkingen"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    examples = [
        "x = 5",
        "y = 3",
        "print(x == y)  # False",
        "print(x != y)  # True",
        "print(x > y)   # True",
        "print(x < y)   # False"
    ]

    for line in examples:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0

    # Slide 4: Meerdere Voorwaarden Combineren
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Meerdere Voorwaarden Combineren"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Gebruik `and`, `or`, en `not` om voorwaarden te combineren:"
    p.level = 0

    combinations = [
        "`and` : Beide voorwaarden moeten waar zijn",
        "`or` : Minstens één voorwaarde moet waar zijn",
        "`not` : Keert de waarheid van de voorwaarde om"
    ]

    for line in combinations:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 5: Voorbeeld met Meerdere Voorwaarden
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Voorbeeld met Meerdere Voorwaarden"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    example_code = [
        "leeftijd = 16",
        "heeft_toestemming = True",
        "",
        "if leeftijd >= 16 and heeft_toestemming:",
        "    print('Je mag met de bromfiets rijden!')",
        "else:",
        "    print('Je mag nog niet met de bromfiets rijden.')"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.strip() else 1

    # Save Presentation
    prs.save('Hoofdstuk_3_1.pptx')
    print("Presentatie 'Hoofdstuk_3_1.pptx' is gegenereerd.")

def create_chapter_3_2_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hoofdstuk 3.2\nIf-statements"
    slide.placeholders[1].text = "Laat je programma beslissingen nemen met if-statements"

    # Slide 2: Het if-statement
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Het if-statement"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Syntax:"
    p.level = 0

    code = [
        "if voorwaarde:",
        "    # code die wordt uitgevoerd als de voorwaarde waar is"
    ]

    for line in code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 3: Voorbeeld van een if-statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Voorbeeld van een if-statement"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    example_code = [
        "leeftijd = 18",
        "",
        "if leeftijd >= 18:",
        "    print('Je bent volwassen.')"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.strip() else 1

    # Slide 4: Het if-else statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Het if-else statement"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Syntax:"
    p.level = 0

    code = [
        "if voorwaarde:",
        "    # code als voorwaarde waar is",
        "else:",
        "    # code als voorwaarde niet waar is"
    ]

    for line in code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 5: Voorbeeld van een if-else statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Voorbeeld van een if-else statement"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    example_code = [
        "punten = 75",
        "",
        "if punten >= 50:",
        "    print('Gefeliciteerd, je bent geslaagd!')",
        "else:",
        "    print('Helaas, je bent niet geslaagd.')"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.strip() else 1

    # Slide 6: Tips voor if-statements
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Tips voor if-statements"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    tips = [
        "1. Inspringing is belangrijk: gebruik consistente indentatie.",
        "2. Vergeet de dubbele punt ':' na de voorwaarde niet.",
        "3. Houd het leesbaar: splits complexe voorwaarden op.",
        "4. Test je voorwaarden met verschillende inputs."
    ]

    for tip in tips:
        p = tf.add_paragraph()
        p.text = tip
        p.level = 1

    # Save Presentation
    prs.save('Hoofdstuk_3_2.pptx')
    print("Presentatie 'Hoofdstuk_3_2.pptx' is gegenereerd.")

if __name__ == "__main__":
    create_chapter_3_1_presentation()
    create_chapter_3_2_presentation()
