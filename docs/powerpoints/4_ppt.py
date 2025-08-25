import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_chapter_4_4_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hoofdstuk 4.4\nDictionaries (Woordenboeken)"
    slide.placeholders[1].text = "Leer hoe je dictionaries gebruikt in Python"

    # Slide 2: Wat zijn Dictionaries?
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Wat zijn Dictionaries?"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Een dictionary is een verzameling van sleutel-waarde paren:"
    p.level = 0

    points = [
        "Ongeordend (tot Python 3.7)",
        "Veranderlijk (Mutable)",
        "Sleutels moeten uniek zijn",
        "Sleutels moeten onveranderlijk (immutable) zijn"
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1

    # Slide 3: Een Dictionary Aanmaken
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Een Dictionary Aanmaken"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    example_code = [
        "# Een lege dictionary",
        "leeg_woordenboek = {}",
        "",
        "# Een dictionary met sleutel-waarde paren",
        "student = {",
        "    'naam': 'Friso',",
        "    'leeftijd': 17,",
        "    'klas': '5BCW',",
        "    'vakken': ['Wiskunde', 'Informatica', 'Nederlands']",
        "}"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.startswith('#') else 1

    # Slide 4: Toegang tot Waarden
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Toegang tot Waarden"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    example_code = [
        "student = {'naam': 'Lobna', 'leeftijd': 17}",
        "",
        "# Direct toegang",
        "print(student['naam'])  # Output: Lobna",
        "",
        "# Veilige toegang met get()",
        "print(student.get('adres', 'Onbekend'))"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.startswith('#') else 1

    # Slide 5: Wijzigen van een Dictionary
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Wijzigen van een Dictionary"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    operations = [
        "Toevoegen: dict['nieuwe_sleutel'] = waarde",
        "Wijzigen: dict['bestaande_sleutel'] = nieuwe_waarde",
        "Verwijderen: del dict['sleutel']",
        "Verwijderen met return: waarde = dict.pop('sleutel')"
    ]

    for op in operations:
        p = tf.add_paragraph()
        p.text = op
        p.level = 1

    # Save Presentation
    prs.save('Hoofdstuk_4_4.pptx')
    print("Presentatie 'Hoofdstuk_4_4.pptx' is gegenereerd.")

def create_chapter_4_5_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hoofdstuk 4.5\nWhile Loops"
    slide.placeholders[1].text = "Leer hoe je while-loops gebruikt in Python"

    # Slide 2: Wat is een While Loop?
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Wat is een While Loop?"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Een while-loop voert code herhaaldelijk uit zolang een voorwaarde waar is:"
    p.level = 0

    steps = [
        "1. Evalueer de voorwaarde",
        "2. Voer code uit als voorwaarde waar is",
        "3. Herhaal tot voorwaarde onwaar wordt"
    ]

    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1

    # Slide 3: Eenvoudige While Loop
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Eenvoudige While Loop"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    example_code = [
        "teller = 1",
        "while teller <= 5:",
        "    print(teller)",
        "    teller += 1  # Verhoog de teller met 1"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0

    # Slide 4: Break en Continue
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Break en Continue Statements"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    statements = [
        "break: Verlaat de loop onmiddellijk",
        "continue: Ga door naar de volgende iteratie",
        "",
        "Voorbeeld:",
        "while True:",
        "    if conditie:",
        "        break    # Verlaat de loop",
        "    if andere_conditie:",
        "        continue # Ga naar volgende iteratie"
    ]

    for line in statements:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line in ["break:", "continue:"] else 1

    # Slide 5: Praktische Toepassingen
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Praktische Toepassingen"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    applications = [
        "1. Gebruikersinvoer valideren",
        "2. Spellen en simulaties",
        "3. Dataverwerking",
        "4. Menu-systemen",
        "5. Wachten op bepaalde condities"
    ]

    for app in applications:
        p = tf.add_paragraph()
        p.text = app
        p.level = 1

    # Save Presentation
    prs.save('Hoofdstuk_4_5.pptx')
    print("Presentatie 'Hoofdstuk_4_5.pptx' is gegenereerd.")

if __name__ == "__main__":
    create_chapter_4_4_presentation()
    create_chapter_4_5_presentation() 