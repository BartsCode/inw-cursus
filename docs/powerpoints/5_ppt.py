import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_chapter_5_1_presentation():
    """
    Generates a PowerPoint presentation for Chapter 5.1: Functies.
    """
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hoofdstuk 5.1\nFuncties"
    slide.placeholders[1].text = "Organiseer en hergebruik je code met functies"

    # Slide 2: Wat zijn Functies?
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Wat zijn Functies?"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Een functie is een blok code met een naam dat een specifieke taak uitvoert."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Voordelen:"
    p.level = 0

    benefits = [
        "Modulair: Code opdelen in kleinere taken.",
        "Herbruikbaar: Dezelfde logica makkelijk opnieuw gebruiken (DRY).",
        "Leesbaarder: Duidelijke namen maken code begrijpelijker.",
        "Onderhoudbaar: Aanpassingen op één centrale plek."
    ]

    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.level = 1

    # Slide 3: Een Functie Definiëren
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Een Functie Definiëren"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Gebruik het `def` sleutelwoord:"
    p.level = 0

    code = [
        "def functienaam(parameters):",
        "    # Code blok (ingesprongen)",
        "    # ... doet iets ...",
        "    return waarde # Optioneel"
    ]

    for line in code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1

    p = tf.add_paragraph()
    p.text = "Voorbeeld:"
    p.level = 0

    example_code = [
        "def groet():",
        "  print(\"Hallo daar!\")"
    ]
    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 4: Een Functie Aanroepen
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Een Functie Aanroepen"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Roep een functie aan met haar naam gevolgd door haakjes `()`:"
    p.level = 0

    example_code = [
        "def groet():",
        "  print(\"Hallo daar!\")",
        "",
        "# Roep de functie aan",
        "groet()",
        "groet() # Kan meerdere keren"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith(" ") or line.startswith("#") or not line.strip() else 0


    # Slide 5: Parameters en Argumenten
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Parameters en Argumenten"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Functies kunnen gegevens ontvangen:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Parameters: Variabelen in de functie definitie `def func(param):`"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Argumenten: Waarden meegegeven bij het aanroepen `func(arg)`"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Voorbeeld:"
    p.level = 0

    example_code = [
        "def begroet_persoon(naam): # 'naam' is parameter",
        "  print(f\"Hallo, {naam}!\")",
        "",
        "begroet_persoon(\"Alice\") # \"Alice\" is argument",
        "begroet_persoon(\"Bob\")"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith(" ") or line.startswith("#") or not line.strip() else 0

    # Slide 6: Return Waarden
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Return Waarden"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Functies kunnen een resultaat teruggeven met `return`."
    p.level = 0

    example_code = [
        "def tel_op(getal1, getal2):",
        "  som = getal1 + getal2",
        "  return som # Geeft de som terug",
        "",
        "resultaat = tel_op(10, 5)",
        "print(f\"Het resultaat is: {resultaat}\")",
        "",
        "# Meerdere waarden teruggeven (als tuple)",
        "def bereken(a, b):",
        "    return a+b, a*b",
        "",
        "s, p = bereken(5, 3)",
        "print(f\"Som: {s}, Product: {p}\")"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith(" ") or line.startswith("#") or not line.strip() else 0

    # Slide 7: Scope van Variabelen
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Scope van Variabelen"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Lokaal: Gedefinieerd binnen een functie, alleen daar geldig."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Globaal: Gedefinieerd buiten functies, overal toegankelijk."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Voorbeeld:"
    p.level = 0

    example_code = [
        "globale_var = \"Ik ben globaal\"",
        "",
        "def mijn_functie():",
        "  lokale_var = \"Ik ben lokaal\"",
        "  print(lokale_var)",
        "  print(globale_var)",
        "",
        "mijn_functie()",
        "# print(lokale_var) # Geeft FOUT!"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith(" ") or line.startswith("#") or not line.strip() else 0

    # Slide 8: Docstrings
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Docstrings"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = "Documenteer je functies met een docstring (`\"\"\"...\"\"\"`):"
    p.level = 0

    example_code = [
        "def tel_op(getal1, getal2):",
        "  \"\"\"",
        "  Telt twee getallen bij elkaar op.",
        "  Args:",
        "    getal1: Het eerste getal.",
        "    getal2: Het tweede getal.",
        "  Returns:",
        "    De som van de twee getallen.",
        "  \"\"\"",
        "  return getal1 + getal2",
        "",
        "# help(tel_op) # Toont de docstring"
    ]

    for line in example_code:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith(" ") or line.startswith("#") or not line.strip() else 0

    # Save Presentation
    try:
        # Ensure the directory exists
        if not os.path.exists('powerpoints'):
            os.makedirs('powerpoints')
        
        output_filename = os.path.join('powerpoints', 'Hoofdstuk_5_1.pptx')
        prs.save(output_filename)
        print(f"Presentatie '{output_filename}' is gegenereerd.")
    except Exception as e:
        print(f"Kon presentatie niet opslaan: {e}")


if __name__ == "__main__":
    create_chapter_5_1_presentation()
